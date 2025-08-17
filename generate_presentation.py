from pptx import Presentation
from pptx.util import Inches

# Create presentation
prs = Presentation()

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Pollution Prevention Overview"
slide.placeholders[1].text = "Webinar for Industrial Process P2"
slide.notes_slide.notes_text_frame.text = (
    "Introduction to pollution prevention and agenda for discussion"
)

# Slide 1: Pollution Prevention Act Overview
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Pollution Prevention Act"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "1990 law prioritizing source reduction"
p = tf.add_paragraph()
p.text = "Waste prevention over treatment and disposal"
p.level = 1
p = tf.add_paragraph()
p.text = "Encourages energy and resource efficiency"
p.level = 1
slide.notes_slide.notes_text_frame.text = (
    "The P2 Act sets a national policy to prevent or reduce pollution at the source (EPA, 1990)."
)

# Slide 2: Tenets of the P2 Act
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Tenets of the P2 Act"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Source reduction is preferred"
p = tf.add_paragraph()
p.text = "Pollution should be prevented or reduced at the source"
p.level = 1
p = tf.add_paragraph()
p.text = "Environmentally safe recycling is next best"
p.level = 1
p = tf.add_paragraph()
p.text = "Treatment is preferable to disposal"
p.level = 1
slide.notes_slide.notes_text_frame.text = (
    "Hierarchy: prevent, recycle, treat, dispose. This guides industrial decision-making (EPA, 1990)."
)

# Slide 3: Clean Water Act
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Clean Water Act"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Regulates discharges into U.S. waters"
p = tf.add_paragraph()
p.text = "NPDES permits control effluent"
p.level = 1
p = tf.add_paragraph()
p.text = "Promotes pollution prevention plans"
p.level = 1
slide.notes_slide.notes_text_frame.text = (
    "Facilities must manage wastewater through permits and adopt BMPs to minimize pollutants."
)

# Slide 4: Clean Air Act
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Clean Air Act"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Establishes National Ambient Air Quality Standards"
p = tf.add_paragraph()
p.text = "Requires control technologies for emission sources"
p.level = 1
p = tf.add_paragraph()
p.text = "Title V operating permits include P2 provisions"
p.level = 1
slide.notes_slide.notes_text_frame.text = (
    "The CAA drives emission reductions and encourages source control and energy efficiency."
)

# Slide 5: Resource Conservation and Recovery Act
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Waste Laws"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "RCRA governs hazardous waste management"
p = tf.add_paragraph()
p.text = "Encourages waste minimization and recycling"
p.level = 1
p = tf.add_paragraph()
p.text = "Facilities track cradle-to-grave wastes"
p.level = 1
slide.notes_slide.notes_text_frame.text = (
    "Subtitle C requires generators to manage waste from generation to disposal and promotes minimization."
)

# Slide 6: Audit Phase 1 & 2
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "P2 Audit Phases 1-2"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Planning and organization"
p = tf.add_paragraph()
p.text = "Pre-audit data gathering"
p.level = 1
p = tf.add_paragraph()
p.text = "On-site assessment of processes"
p.level = 1
slide.notes_slide.notes_text_frame.text = (
    "Phase 1 defines scope and teams; Phase 2 maps material and energy flows."
)

# Slide 7: Audit Phase 3 & 4
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "P2 Audit Phases 3-4"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Feasibility analysis of options"
p = tf.add_paragraph()
p.text = "Implementation and monitoring"
p.level = 1
slide.notes_slide.notes_text_frame.text = (
    "Options are prioritized by cost and impact; progress is tracked after implementation."
)

# Slide 8: Assessing Pollutant Source Control
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Assessing Source Control"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Material balance calculations"
p = tf.add_paragraph()
p.text = "Process mapping and benchmarking"
p.level = 1
p = tf.add_paragraph()
p.text = "Identifies root causes of waste"
p.level = 1
slide.notes_slide.notes_text_frame.text = (
    "Audits quantify inputs/outputs to locate inefficiencies and prioritize source reduction."
)

# Slide 9: Life Cycle Analysis
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Life Cycle Analysis"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Evaluates environmental impacts across stages"
p = tf.add_paragraph()
p.text = "Considers raw materials, production, use, disposal"
p.level = 1
slide.notes_slide.notes_text_frame.text = (
    "LCA helps compare alternatives and avoid burden shifting (Sikdar, 2003)."
)

# Slide 10: Cradle-to-Grave Value
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Cradle-to-Grave Management"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Holistic view of material stewardship"
p = tf.add_paragraph()
p.text = "Reveals opportunities for recycling and design changes"
p.level = 1
p = tf.add_paragraph()
p.text = "Supports sustainable decision-making"
p.level = 1
slide.notes_slide.notes_text_frame.text = (
    "Understanding life cycle ensures responsibilities for impacts at each stage are addressed."
)

# Slide 11: Conclusion
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Key Takeaways"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "P2 Act promotes source reduction"
p = tf.add_paragraph()
p.text = "Federal laws integrate P2 in water, air, and waste programs"
p.level = 1
p = tf.add_paragraph()
p.text = "Audits and LCA guide sustainable practices"
p.level = 1
slide.notes_slide.notes_text_frame.text = (
    "Summarizes benefits of integrating P2 and LCA in industrial operations."
)

# References slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "References"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Environmental Protection Agency. (1990). Pollution Prevention Act."
p = tf.add_paragraph()
p.text = "Sikdar, S. K. (2003). Sustainable development and pollution prevention. Journal of Cleaner Production, 11(1), 131-145."
p.level = 1
slide.notes_slide.notes_text_frame.text = (
    "Both references provide foundational information on P2 and sustainable development."
)

prs.save("pollution_prevention_presentation.pptx")
